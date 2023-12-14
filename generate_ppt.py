from pptx import Presentation
from pptx.util import Inches


def create_presentation():
    prs = Presentation()

    # Slide 1: What is old agriculture and its drawbacks?
    slide_1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_1.shapes.title.text = "What is old agriculture and its drawbacks?"
    content = (
        "- Definition: Old agriculture, often termed as traditional farming, involves farming methods used before the industrial revolution.\n"
        "- Drawbacks:\n"
        "  - Limited Production: Dependent on natural conditions.\n"
        "  - Labor-Intensive: Required a lot of manual work.\n"
        "  - Unsustainable: Led to soil erosion and degradation.\n"
        "  - Limited Tools: Dependent on basic tools and animals."
    )
    slide_1.placeholders[1].text = content

    # Slide 2: Why we need modern agriculture?
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_2.shapes.title.text = "Why we need modern agriculture?"
    content = (
        "- Higher Demand: Rising global population.\n"
        "- Sustainability: Need for eco-friendly methods.\n"
        "- Efficiency: Modern tools increase yield with less effort.\n"
        "- Diversity: Modern methods allow for a broader range of crops."
    )
    slide_2.placeholders[1].text = content

    # ... (similarly for other slides)

    # Slide 7: Upcoming modern agriculture technologies
    slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_7.shapes.title.text = "Upcoming modern agriculture technologies"
    content = (
        "- AI and Machine Learning for crop predictions.\n"
        "- Drone technology for monitoring fields.\n"
        "- CRISPR technology for plant breeding.\n"
        "- Autonomous tractors and robots for farming."
    )
    slide_7.placeholders[1].text = content

    prs.save("Agriculture_Presentation.pptx")


if __name__ == "__main__":
    create_presentation()
